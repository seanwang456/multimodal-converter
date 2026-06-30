"""Phase 4 复杂文档互转验证：pdf→txt/docx/xlsx、docx↔pptx。"""
from __future__ import annotations

import asyncio
from pathlib import Path

from app.handlers.docgen import write_pdf
from app.handlers.office_handlers import OFFICE_HANDLERS
from app.handlers.pdf_handlers import PDF_HANDLERS


def _run(handler, input_path: Path, out_dir: Path, ext: str) -> dict:
    return asyncio.run(handler.run(str(input_path), str(out_dir), ext, {}))


def _make_pdf(p: Path) -> Path:
    write_pdf("你好 PDF 内容\n第二行 文字", p)
    return p


def _make_docx(p: Path) -> Path:
    from docx import Document
    d = Document(); d.add_paragraph("标题段"); d.add_paragraph("正文内容 ABC")
    d.save(str(p)); return p


def _make_pptx(p: Path) -> Path:
    from pptx import Presentation
    prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "幻灯标题"; s.placeholders[1].text = "正文要点"
    prs.save(str(p)); return p


def test_pdf_to_txt(tmp_path) -> None:
    src = _make_pdf(tmp_path / "s.pdf")
    out_dir = tmp_path / "out"; out_dir.mkdir()
    res = _run(PDF_HANDLERS["pdf_to_txt"], src, out_dir, ".txt")
    content = Path(res["output_path"]).read_text(encoding="utf-8")
    assert "PDF" in content  # 中文提取随库，至少英文/关键字应存在


def test_pdf_to_docx(tmp_path) -> None:
    src = _make_pdf(tmp_path / "s.pdf")
    out_dir = tmp_path / "out"; out_dir.mkdir()
    res = _run(PDF_HANDLERS["pdf_to_docx"], src, out_dir, ".docx")
    out = Path(res["output_path"])
    from docx import Document
    Document(str(out))  # 可打开
    assert out.suffix == ".docx" and out.stat().st_size > 0


def test_pdf_to_xlsx_no_table(tmp_path) -> None:
    from app.errors import ConversionError, ErrorCode
    src = _make_pdf(tmp_path / "s.pdf")  # 纯文本 PDF，无表格
    out_dir = tmp_path / "out"; out_dir.mkdir()
    try:
        _run(PDF_HANDLERS["pdf_to_xlsx"], src, out_dir, ".xlsx")
        assert False, "应抛 NO_TABLE_FOUND"
    except ConversionError as e:
        assert e.code == ErrorCode.NO_TABLE_FOUND


def test_docx_to_pptx(tmp_path) -> None:
    src = _make_docx(tmp_path / "s.docx")
    out_dir = tmp_path / "out"; out_dir.mkdir()
    res = _run(OFFICE_HANDLERS["docx_to_pptx"], src, out_dir, ".pptx")
    out = Path(res["output_path"])
    from pptx import Presentation
    Presentation(str(out))
    assert out.suffix == ".pptx" and out.stat().st_size > 0


def test_pptx_to_docx(tmp_path) -> None:
    src = _make_pptx(tmp_path / "s.pptx")
    out_dir = tmp_path / "out"; out_dir.mkdir()
    res = _run(OFFICE_HANDLERS["pptx_to_docx"], src, out_dir, ".docx")
    out = Path(res["output_path"])
    from docx import Document
    Document(str(out))
    assert out.suffix == ".docx" and out.stat().st_size > 0


def test_registry_phase4_routes() -> None:
    from app.services import registry
    registry.register_handlers()
    pdf_targets = {r["target_ext"] for r in registry.get_targets(".pdf")}
    assert {".txt", ".docx", ".xlsx", ".pptx"} <= pdf_targets
    docx_targets = {r["target_ext"] for r in registry.get_targets(".docx")}
    assert ".pptx" in docx_targets
    pptx_targets = {r["target_ext"] for r in registry.get_targets(".pptx")}
    assert ".docx" in pptx_targets
