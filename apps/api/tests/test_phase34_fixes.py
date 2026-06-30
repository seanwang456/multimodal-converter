"""Codex Phase 3/4 REVISE 修复的回归测试。"""
from __future__ import annotations

import asyncio
from pathlib import Path
from types import SimpleNamespace

import pytest

from app.errors import ConversionError, ErrorCode
from app.handlers.multimodal_handlers import MULTIMODAL_HANDLERS
from app.handlers.office_handlers import OFFICE_HANDLERS
from app.handlers.pdf_handlers import PDF_HANDLERS


def _run(handler, input_path: Path, out_dir: Path, ext: str) -> dict:
    return asyncio.run(handler.run(str(input_path), str(out_dir), ext, {}))


def _no_creds():
    return SimpleNamespace(llm_base_url="", llm_api_key="")


# ---------- bmp OCR 归一化 ----------
def test_image_data_url_normalizes_bmp(tmp_path) -> None:
    from PIL import Image
    from app.providers.llm_provider import _image_data_url

    p = tmp_path / "b.bmp"
    Image.new("RGB", (4, 4), (1, 2, 3)).save(str(p), format="BMP")
    url = _image_data_url(str(p))
    assert url.startswith("data:image/png;base64,")


# ---------- 缺凭证错误码 ----------
@pytest.mark.asyncio
async def test_ocr_missing_creds(monkeypatch) -> None:
    from app.providers import llm_provider

    monkeypatch.setattr(llm_provider, "settings", _no_creds())
    # _creds 会回退到 os.getenv；.env 经 load_dotenv 已注入环境，需一并清除才真正"无凭证"
    monkeypatch.delenv("LLM_BASE_URL", raising=False)
    monkeypatch.delenv("LLM_API_KEY", raising=False)
    try:
        await llm_provider.OpenAIOCRProvider().extract_text("x.png")
        assert False
    except ConversionError as e:
        assert e.code == ErrorCode.OCR_FAILED


@pytest.mark.asyncio
async def test_asr_missing_creds(monkeypatch) -> None:
    from app.providers import llm_provider

    monkeypatch.setattr(llm_provider, "settings", _no_creds())
    # _creds 会回退到 os.getenv；.env 经 load_dotenv 已注入环境，需一并清除才真正"无凭证"
    monkeypatch.delenv("LLM_BASE_URL", raising=False)
    monkeypatch.delenv("LLM_API_KEY", raising=False)
    try:
        await llm_provider.OpenAIASRProvider().transcribe("x.wav")
        assert False
    except ConversionError as e:
        assert e.code == ErrorCode.ASR_FAILED


# ---------- 加密 PDF 全目标拒绝 ----------
def _make_encrypted_pdf(p: Path) -> Path:
    from reportlab.pdfgen import canvas
    from pypdf import PdfReader, PdfWriter

    canvas.Canvas(str(p)).save()  # 空白 PDF
    writer = PdfWriter(clone_from=str(p))
    writer.encrypt("secret")
    with open(p, "wb") as f:
        writer.write(f)
    return p


@pytest.mark.parametrize("key,target", [
    ("pdf_to_txt", ".txt"), ("pdf_to_docx", ".docx"),
    ("pdf_to_pptx", ".pptx"), ("pdf_to_xlsx", ".xlsx"),
])
def test_encrypted_pdf_rejected_all_targets(tmp_path, key, target) -> None:
    src = _make_encrypted_pdf(tmp_path / "enc.pdf")
    out_dir = tmp_path / "out"; out_dir.mkdir()
    try:
        _run(PDF_HANDLERS[key], src, out_dir, target)
        assert False, "加密 PDF 应被拒绝"
    except ConversionError as e:
        assert e.code == ErrorCode.PASSWORD_PROTECTED_PDF


# ---------- docx↔pptx 内容保留 ----------
def test_docx_to_pptx_preserves_text(tmp_path) -> None:
    from docx import Document

    src = tmp_path / "s.docx"
    d = Document(); d.add_paragraph("独特标题XYZ"); d.add_paragraph("第二段正文ABC")
    d.save(str(src))
    out_dir = tmp_path / "out"; out_dir.mkdir()
    res = _run(OFFICE_HANDLERS["docx_to_pptx"], src, out_dir, ".pptx")
    from pptx import Presentation
    texts = []
    for slide in Presentation(str(res["output_path"])).slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text_frame.text)
    joined = "\n".join(texts)
    assert "独特标题XYZ" in joined


def test_pptx_to_docx_preserves_text(tmp_path) -> None:
    from pptx import Presentation

    src = tmp_path / "s.pptx"
    prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "幻灯标题999"; s.placeholders[1].text = "正文要点888"
    prs.save(str(src))
    out_dir = tmp_path / "out"; out_dir.mkdir()
    res = _run(OFFICE_HANDLERS["pptx_to_docx"], src, out_dir, ".docx")
    from docx import Document
    joined = "\n".join(p.text for p in Document(str(res["output_path"])).paragraphs)
    assert "幻灯标题999" in joined and "正文要点888" in joined
