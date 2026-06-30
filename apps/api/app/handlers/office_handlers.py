"""Office 相关：office_to_pdf(LibreOffice 异步子进程)、office_to_txt/xlsx_to_csv(库提取，to_thread)。规格 §7.2.1。"""
from __future__ import annotations

import asyncio
import csv
from pathlib import Path
from typing import Any

from app.config import settings
from app.errors import ConversionError, ErrorCode
from app.handlers._subprocess import run_subprocess
from app.handlers.base import ConversionHandler, ConversionResult

OFFICE_SRCS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}


def _extract_text(input_path: str) -> str:
    src = Path(input_path).suffix.lower()
    if src in (".docx", ".doc"):
        from docx import Document
        return "\n".join(p.text for p in Document(input_path).paragraphs)
    if src in (".pptx", ".ppt"):
        from pptx import Presentation
        lines: list[str] = []
        for slide in Presentation(input_path).slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)
    if src in (".xlsx", ".xls"):
        from openpyxl import load_workbook
        wb = load_workbook(input_path, read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join("" if c is None else str(c) for c in row))
        return "\n".join(rows)
    raise ConversionError(ErrorCode.UNSUPPORTED_CONVERSION, f"不支持 {src} 提取文本")


class OfficeToPdfHandler(ConversionHandler):
    """docx/pptx/xlsx/doc/ppt/xls → pdf，LibreOffice headless（每 job 独立 profile 防并发锁死）。"""

    key = "office_to_pdf"

    def supports(self, source_ext: str, target_ext: str) -> bool:
        return source_ext in OFFICE_SRCS and target_ext == ".pdf"

    async def run(
        self, input_path: str, output_dir: str, target_ext: str, options: dict[str, Any]
    ) -> ConversionResult:
        out_dir = Path(output_dir)
        out = out_dir / (Path(input_path).stem + ".pdf")
        profile = (out_dir / "lo_profile").as_uri()
        cmd = [
            settings.libreoffice_bin,
            f"-env:UserInstallation={profile}",
            "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), input_path,
        ]
        try:
            rc, _out, _err = await run_subprocess(cmd, 180)
        except asyncio.TimeoutError as e:
            raise ConversionError(ErrorCode.CONVERSION_TIMEOUT, "文档转换超时") from e
        except FileNotFoundError as e:
            raise ConversionError(ErrorCode.CONVERSION_ENGINE_ERROR, "LibreOffice 未安装") from e
        if rc != 0 or not out.exists():
            raise ConversionError(ErrorCode.CONVERSION_ENGINE_ERROR, "文档转换失败")
        return ConversionResult(
            output_path=str(out), filename=out.name, size_bytes=out.stat().st_size,
            quality_notice="文档转 PDF 尽量保持原样，复杂版式可能略有差异。",
        )


class OfficeToTxtHandler(ConversionHandler):
    key = "office_to_txt"

    def supports(self, source_ext: str, target_ext: str) -> bool:
        return source_ext in OFFICE_SRCS and target_ext == ".txt"

    async def run(
        self, input_path: str, output_dir: str, target_ext: str, options: dict[str, Any]
    ) -> ConversionResult:
        text = await asyncio.to_thread(_extract_text, input_path)
        out = Path(output_dir) / "result.txt"
        await asyncio.to_thread(out.write_text, text, "utf-8")
        return ConversionResult(
            output_path=str(out), filename=out.name, size_bytes=out.stat().st_size, quality_notice=None
        )


class XlsxToCsvHandler(ConversionHandler):
    key = "xlsx_to_csv"

    def supports(self, source_ext: str, target_ext: str) -> bool:
        return source_ext in {".xlsx", ".xls"} and target_ext == ".csv"

    async def run(
        self, input_path: str, output_dir: str, target_ext: str, options: dict[str, Any]
    ) -> ConversionResult:
        out = Path(output_dir) / "result.csv"
        await asyncio.to_thread(_write_csv, input_path, out)
        return ConversionResult(
            output_path=str(out), filename=out.name, size_bytes=out.stat().st_size,
            quality_notice="仅导出第一个工作表。",
        )


def _write_csv(input_path: str, out: Path) -> None:
    from openpyxl import load_workbook

    wb = load_workbook(input_path, read_only=True, data_only=True)
    ws = wb.worksheets[0]  # 仅第一个 sheet
    with open(out, "w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        for row in ws.iter_rows(values_only=True):
            writer.writerow(["" if c is None else c for c in row])


class DocxToPptxHandler(ConversionHandler):
    """docx → pptx：按段落生成幻灯片。规格 §7.2.1。"""

    key = "docx_to_pptx"

    def supports(self, source_ext: str, target_ext: str) -> bool:
        return source_ext in {".docx", ".doc"} and target_ext == ".pptx"

    async def run(self, input_path: str, output_dir: str, target_ext: str, options: dict[str, Any]) -> ConversionResult:
        from pptx import Presentation
        from pptx.util import Inches

        def _build() -> Path:
            from docx import Document

            paras = [p.text.strip() for p in Document(input_path).paragraphs if p.text.strip()]
            prs = Presentation()
            layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
            for para in paras or [""]:
                slide = prs.slides.add_slide(layout)
                title_ph = getattr(slide.shapes, "title", None)
                if title_ph is not None:
                    title_ph.text = para[:60]
                # 仅当段落超出标题长度时才补正文框，避免短段落标题与正文重复
                if len(para) > 60:
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                    box.text_frame.word_wrap = True
                    box.text_frame.text = para
            out = Path(output_dir) / "result.pptx"
            prs.save(str(out))
            return out

        out = await asyncio.to_thread(_build)
        return ConversionResult(
            output_path=str(out), filename=out.name, size_bytes=out.stat().st_size,
            quality_notice="Word 转 PPT 按段落分页，丢失复杂版式。",
        )


class PptxToDocxHandler(ConversionHandler):
    """pptx → docx：每页文字转段落。规格 §7.2.1。"""

    key = "pptx_to_docx"

    def supports(self, source_ext: str, target_ext: str) -> bool:
        return source_ext in {".pptx", ".ppt"} and target_ext == ".docx"

    async def run(self, input_path: str, output_dir: str, target_ext: str, options: dict[str, Any]) -> ConversionResult:
        def _build() -> Path:
            from docx import Document
            from pptx import Presentation

            doc = Document()
            for slide in Presentation(input_path).slides:
                texts = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                        texts.append(shape.text_frame.text.strip())
                if texts:
                    doc.add_heading(texts[0][:60], level=2)
                    for t in texts[1:]:
                        doc.add_paragraph(t)
            out = Path(output_dir) / "result.docx"
            doc.save(str(out))
            return out

        out = await asyncio.to_thread(_build)
        return ConversionResult(
            output_path=str(out), filename=out.name, size_bytes=out.stat().st_size,
            quality_notice="PPT 转 Word 仅提取文字，丢失版式与图片。",
        )


OFFICE_HANDLERS: dict[str, ConversionHandler] = {
    "office_to_pdf": OfficeToPdfHandler(),
    "office_to_txt": OfficeToTxtHandler(),
    "xlsx_to_csv": XlsxToCsvHandler(),
    "docx_to_pptx": DocxToPptxHandler(),
    "pptx_to_docx": PptxToDocxHandler(),
}
