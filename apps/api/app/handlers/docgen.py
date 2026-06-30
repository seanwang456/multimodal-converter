"""文本 → docx/pdf/pptx/xlsx 生成（txt handler 与 OCR/ASR handler 共享）。"""
from __future__ import annotations

from pathlib import Path


def write_docx(text: str, dest: Path) -> Path:
    from docx import Document

    doc = Document()
    for line in text.splitlines() or [""]:
        doc.add_paragraph(line)
    doc.save(str(dest))
    return dest


def write_pdf(text: str, dest: Path) -> Path:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate
    from xml.sax.saxutils import escape

    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))  # 支持中文
        font = "STSong-Light"
    except Exception:  # noqa: BLE001
        font = "Helvetica"

    base = getSampleStyleSheet()["Normal"]
    style = ParagraphStyle("body", parent=base, fontName=font, fontSize=11, leading=16)
    doc = SimpleDocTemplate(str(dest), pagesize=A4,
                            leftMargin=40, rightMargin=40, topMargin=40, bottomMargin=40)
    flow = []
    for line in text.splitlines():
        # Platypus 会按页宽自动换行，长行不再截断
        flow.append(Paragraph(escape(line) if line.strip() else "&nbsp;", style))
    if not flow:
        flow.append(Paragraph("&nbsp;", style))
    doc.build(flow)
    return dest


def write_pptx(text: str, dest: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
    for page in (text.split("\n\n") or [text]):
        slide = prs.slides.add_slide(layout)
        lines = [ln for ln in page.splitlines() if ln.strip()]
        if not lines:
            continue
        title_ph = getattr(slide.shapes, "title", None)
        if title_ph is not None:
            title_ph.text = lines[0][:80]
        body = "\n".join(lines[1:]) if len(lines) > 1 else lines[0]
        # 用 Inches 让文本框尺寸真实可见（裸整数会被当作 EMU ≈ 0）
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = body
    prs.save(str(dest))
    return dest


def write_xlsx(text: str, dest: Path) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "sheet1"
    for row_idx, line in enumerate(text.splitlines(), start=1):
        cells = line.split("\t") if "\t" in line else (line.split(",") if "," in line else [line])
        for col_idx, val in enumerate(cells, start=1):
            ws.cell(row=row_idx, column=col_idx, value=val)
    wb.save(str(dest))
    return dest
